import os
import io
from typing import Optional, Dict, Any
import PyPDF2
from docx import Document
from pptx import Presentation
import logging
from app.config.settings import settings

logger = logging.getLogger(__name__)

class FileProcessor:
    def __init__(self):
        self.allowed_extensions = settings.ALLOWED_EXTENSIONS
        self.max_file_size = settings.MAX_FILE_SIZE
    
    def validate_file(self, filename: str, file_size: int) -> tuple[bool, str]:
        """Validate file extension and size"""
        if file_size > self.max_file_size:
            return False, f"File size exceeds maximum limit of {self.max_file_size / 1024 / 1024:.1f}MB"
        
        file_extension = filename.lower().split('.')[-1]
        if file_extension not in self.allowed_extensions:
            return False, f"File type '{file_extension}' not supported. Allowed types: {', '.join(self.allowed_extensions)}"
        
        return True, "Valid file"
    
    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file"""
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text = ""
            
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"PDF text extraction failed: {e}")
            return ""
    
    def extract_text_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(io.BytesIO(file_content))
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"DOCX text extraction failed: {e}")
            return ""
    
    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"PPTX text extraction failed: {e}")
            return ""
    
    def extract_text_from_txt(self, file_content: bytes) -> str:
        """Extract text from TXT file"""
        try:
            return file_content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return file_content.decode('latin-1')
            except Exception as e:
                logger.error(f"TXT text extraction failed: {e}")
                return ""
    
    def extract_text(self, file_content: bytes, filename: str) -> str:
        """Extract text from file based on extension"""
        file_extension = filename.lower().split('.')[-1]
        
        extractors = {
            'pdf': self.extract_text_from_pdf,
            'docx': self.extract_text_from_docx,
            'doc': self.extract_text_from_docx,  # Assuming DOCX format
            'pptx': self.extract_text_from_pptx,
            'ppt': self.extract_text_from_pptx,  # Assuming PPTX format
            'txt': self.extract_text_from_txt
        }
        
        extractor = extractors.get(file_extension)
        if extractor:
            return extractor(file_content)
        else:
            logger.error(f"No extractor found for file type: {file_extension}")
            return ""
    
    def get_file_metadata(self, filename: str, file_size: int, content_type: str) -> Dict[str, Any]:
        """Get file metadata"""
        file_extension = filename.lower().split('.')[-1]
        
        return {
            'original_filename': filename,
            'file_type': file_extension,
            'file_size': file_size,
            'content_type': content_type,
            'supported': file_extension in self.allowed_extensions
        }

# Global file processor instance
file_processor = FileProcessor()
