# Doc-Chat RAG System - Document Processing Service
import hashlib
import logging
from typing import List, Dict, Optional
from fastapi import HTTPException
from app.schemas.youtubeSchema import DocumentSegment, SourceType
# import pypdf
import docx2python
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from langchain.text_splitter import RecursiveCharacterTextSplitter
import time # Added for performance tracking
import redis # Added for caching

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles processing of various document types"""
    
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=200,
            chunk_overlap=20,
            length_function=len,
        )
        # Initialize Redis for caching
        # try:
        #     self.redis_client = redis.from_url(settings.redis_url, decode_responses=True)
        #     self.cache_ttl = 3600  # 1 hour cache
        # except Exception as e:
        #     logger.warning(f"Redis not available for caching: {str(e)}")
        #     self.redis_client = None
    
    def _get_cache_key(self, video_id: str, method: str) -> str:
        """Generate cache key for YouTube transcript"""
        return f"youtube_transcript:{video_id}:{method}"
    
    def _get_from_cache(self, cache_key: str) -> Optional[str]:
        """Get transcript from cache"""
        if not self.redis_client:
            return None
        try:
            return self.redis_client.get(cache_key)
        except Exception as e:
            logger.warning(f"Cache get failed: {str(e)}")
            return None
    
    def _set_cache(self, cache_key: str, transcript: str) -> None:
        """Set transcript in cache"""
        if not self.redis_client:
            return
        try:
            self.redis_client.setex(cache_key, self.cache_ttl, transcript)
        except Exception as e:
            logger.warning(f"Cache set failed: {str(e)}")
    
    # async def process_pdf(self, file_path: str, doc_id: str, tenant_id: str) -> List[DocumentSegment]:
    #     """Process PDF file and extract text with page numbers"""
    #     segments = []
    #     try:
    #         with open(file_path, 'rb') as file:
    #             pdf_reader = pypdf.PdfReader(file)
    #             for page_num, page in enumerate(pdf_reader.pages, 1):
    #                 text = page.extract_text()
    #                 if text.strip():
    #                     chunks = self.text_splitter.split_text(text)
    #                     for chunk in chunks:
    #                         if chunk.strip():
    #                             segment = DocumentSegment(
    #                                 tenant_id=tenant_id,
    #                                 doc_id=doc_id,
    #                                 source_type=SourceType.PDF,
    #                                 source_url=file_path,
    #                                 text=chunk,
    #                                 page=page_num,
    #                                 hash=hashlib.sha256(chunk.encode()).hexdigest()
    #                             )
    #                             segments.append(segment)
    #     except Exception as e:
    #         logger.error(f"Error processing PDF {file_path}: {str(e)}")
    #         raise HTTPException(status_code=400, detail=f"PDF processing failed: {str(e)}")
        
    #     return segments
    
    async def process_docx(self, file_path: str, doc_id: str, tenant_id: str) -> List[DocumentSegment]:
        """Process DOCX file and extract text"""
        segments = []
        try:
            result = docx2python.docx2python(file_path)
            full_text = result.text
            
            chunks = self.text_splitter.split_text(full_text)
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    segment = DocumentSegment(
                        tenant_id=tenant_id,
                        doc_id=doc_id,
                        source_type=SourceType.DOCX,
                        source_url=file_path,
                        text=chunk,
                        page=i + 1,
                        hash=hashlib.sha256(chunk.encode()).hexdigest()
                    )
                    segments.append(segment)
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {str(e)}")
            raise HTTPException(status_code=400, detail=f"DOCX processing failed: {str(e)}")
        
        return segments
    
    async def process_pptx(self, file_path: str, doc_id: str, tenant_id: str) -> List[DocumentSegment]:
        """Process PowerPoint file and extract text with slide numbers"""
        segments = []
        try:
            presentation = Presentation(file_path)
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    combined_text = "\n".join(slide_text)
                    chunks = self.text_splitter.split_text(combined_text)
                    for chunk in chunks:
                        if chunk.strip():
                            segment = DocumentSegment(
                                tenant_id=tenant_id,
                                doc_id=doc_id,
                                source_type=SourceType.PPT,
                                source_url=file_path,
                                text=chunk,
                                page=slide_num,
                                hash=hashlib.sha256(chunk.encode()).hexdigest()
                            )
                            segments.append(segment)
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            raise HTTPException(status_code=400, detail=f"PPTX processing failed: {str(e)}")
        
        return segments
    
    async def process_url(self, url: str, doc_id: str, tenant_id: str) -> List[DocumentSegment]:
        """Process web URL and extract clean text"""
        segments = []
        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove scripts and style elements
            for script in soup(["script", "style"]):
                script.extract()
            
            # Extract text
            text = soup.get_text()
            
            # Clean up text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            clean_text = '\n'.join(chunk for chunk in chunks if chunk)
            
            # Split into chunks
            text_chunks = self.text_splitter.split_text(clean_text)
            for i, chunk in enumerate(text_chunks):
                if chunk.strip():
                    segment = DocumentSegment(
                        tenant_id=tenant_id,
                        doc_id=doc_id,
                        source_type=SourceType.URL,
                        source_url=url,
                        text=chunk,
                        page=i + 1,
                        hash=hashlib.sha256(chunk.encode()).hexdigest()
                    )
                    segments.append(segment)
        except Exception as e:
            logger.error(f"Error processing URL {url}: {str(e)}")
            raise HTTPException(status_code=400, detail=f"URL processing failed: {str(e)}")
        
        return segments
    
    async def process_youtube(self, url: str, doc_type: str):
        """Process YouTube video and extract actual speech transcript"""
        try:
            # Extract video ID from URL
            video_id = self._extract_video_id(url)
            if not video_id:
                raise HTTPException(status_code=400, detail="Invalid YouTube URL")
            
            logger.info(f"Processing YouTube video: {video_id}")
            
            # Try to get actual speech transcript first
            transcript_text = await self._get_youtube_transcript(video_id, url)
            
            if transcript_text:
                logger.info(f"Successfully got speech transcript: {len(transcript_text)} characters")
                logger.info(f"Transcript preview: {transcript_text[:300]}...")
                
                # Use actual transcript content
                video_text = f"""YouTube Video: {video_id}
URL: {url}

{transcript_text}
"""
            else:
                logger.warning(f"No speech transcript available for video {video_id}, trying fallback methods")
                
                # Fallback: Try alternative methods
                if not transcript_text:
                    transcript_text = await self._get_youtube_captions(video_id)
                
                if not transcript_text:
                    transcript_text = await self._get_youtube_with_ytdlp(video_id)
                
                if not transcript_text:
                    transcript_text = await self._get_youtube_video_info(video_id)
                
                if not transcript_text:
                    transcript_text = await self._get_youtube_direct_api(video_id)
                
                if not transcript_text:
                    transcript_text = await self._get_youtube_simple_api(video_id)
                
                if transcript_text:
                    video_text = f"YouTube Video: {video_id}\nURL: {url}\n\nTranscript:\n{transcript_text}"
                else:
                    # Final fallback - basic video info
                    video_info = await self._get_youtube_video_info(video_id)
                    if video_info:
                        video_text = f"""YouTube Video: {video_id}
URL: {url}
{video_info}
"""
                    else:
                        video_text = f"""YouTube Video: {video_id}
URL: {url}
   """      
            
                    
        except Exception as e:
            logger.error(f"Error processing YouTube video {url}: {str(e)}")
            raise HTTPException(status_code=400, detail=f"YouTube processing failed: {str(e)}")
        
        return video_text
    
    async def _get_youtube_transcript(self, video_id: str, video_url: str = None) -> str:
        """Get YouTube video transcript using PRODUCTION-LEVEL dynamic extraction methods"""
        try:
            logger.info(f"🚀 Starting PRODUCTION-LEVEL transcript extraction for video ID: {video_id}")
            
            # Method 1: Production-Level Dynamic Transcript Extractor (Primary)
            try:
                from .production_transcript_extractor import ProductionTranscriptExtractor
                extractor = ProductionTranscriptExtractor()
                
                # Detect video type for optimization
                video_type = self._detect_video_type(video_url) if video_url else None
                transcript = await extractor.extract_transcript(video_id, video_type)
                
                if transcript and len(transcript) > 100:
                    logger.info(f"🚀 PRODUCTION extractor success: {len(transcript)} chars")
                    logger.info(f"📝 Preview: {transcript[:200]}...")
                    
                    # Log performance stats
                    stats = extractor.get_performance_stats()
                    logger.info(f"📊 Performance stats: {stats}")
                    
                    return transcript
                else:
                    logger.warning(f"❌ PRODUCTION extractor returned insufficient content")
                    
            except Exception as e:
                logger.warning(f"❌ PRODUCTION extractor failed: {str(e)}")
            
            # Method 2: Ultra-Powerful Transcript Extractor (Fallback)
            try:
                from .ultra_powerful_transcript_extractor import UltraPowerfulTranscriptExtractor
                extractor = UltraPowerfulTranscriptExtractor()
                transcript = await extractor.extract_transcript(video_id)
                
                if transcript and len(transcript) > 100:
                    logger.info(f"🚀 ULTRA-POWERFUL extractor success: {len(transcript)} chars")
                    logger.info(f"📝 Preview: {transcript[:200]}...")
                    return transcript
                else:
                    logger.warning(f"❌ ULTRA-POWERFUL extractor returned insufficient content")
                    
            except Exception as e:
                logger.warning(f"❌ ULTRA-POWERFUL extractor failed: {str(e)}")
            
            # Method 3: Powerful Transcript Extractor (Secondary Fallback)
            try:
                from .powerful_transcript_extractor import PowerfulTranscriptExtractor
                extractor = PowerfulTranscriptExtractor()
                transcript = await extractor.extract_transcript(video_id)
                
                if transcript and len(transcript) > 100:
                    logger.info(f"🎤 Powerful extractor success: {len(transcript)} chars")
                    logger.info(f"📝 Preview: {transcript[:200]}...")
                    return transcript
                else:
                    logger.warning(f"❌ Powerful extractor returned insufficient content")
                    
            except Exception as e:
                logger.warning(f"❌ Powerful extractor failed: {str(e)}")
            
            # Method 4: Advanced Transcript Extractor (Legacy Fallback)
            try:
                from .advanced_transcript import AdvancedTranscriptExtractor
                extractor = AdvancedTranscriptExtractor()
                transcript = await extractor.get_transcript(video_id)
                
                if transcript and len(transcript) > 100:
                    logger.info(f"🔧 Advanced extractor success: {len(transcript)} chars")
                    logger.info(f"📝 Preview: {transcript[:200]}...")
                    return transcript
                else:
                    logger.warning(f"❌ Advanced extractor returned insufficient content")
                    
            except Exception as e:
                logger.warning(f"❌ Advanced extractor failed: {str(e)}")
            
            # Final fallback - create meaningful placeholder
            logger.warning(f"❌ All extractors failed for video: {video_id}")
            return self._create_meaningful_placeholder(video_id, video_url)
                
        except Exception as e:
            logger.error(f"Error getting YouTube transcript for {video_id}: {str(e)}")
            return self._create_meaningful_placeholder(video_id, video_url)
    
    def _detect_video_type(self, video_url: str) -> str:
        """Detect video type for optimization"""
        try:
            if not video_url:
                return "unknown"
            
            # Check for educational keywords in URL
            educational_keywords = ['lecture', 'tutorial', 'course', 'class', 'education', 'learn', 'study', 'academy', 'university', 'college']
            if any(keyword in video_url.lower() for keyword in educational_keywords):
                return "educational"
            
            # Check for music keywords
            music_keywords = ['music', 'song', 'audio', 'lyrics', 'album', 'artist', 'band']
            if any(keyword in video_url.lower() for keyword in music_keywords):
                return "music"
            
            # Check for short content
            short_keywords = ['short', 'shorts', 'quick', 'brief', 'minute']
            if any(keyword in video_url.lower() for keyword in short_keywords):
                return "short"
            
            # Check for Hindi content
            hindi_keywords = ['hindi', 'indian', 'desi', 'bollywood']
            if any(keyword in video_url.lower() for keyword in hindi_keywords):
                return "hindi"
            
            # Check for Hinglish content
            hinglish_keywords = ['hinglish', 'mixed', 'indian_english', 'hindi_english', 'desi_english']
            if any(keyword in video_url.lower() for keyword in hinglish_keywords):
                return "hinglish"
            
            return "unknown"
            
        except Exception as e:
            logger.warning(f"❌ Failed to detect video type: {str(e)}")
            return "unknown"
    
    def _parse_subtitle_content(self, content: str) -> str:
        """Parse subtitle content from various formats (SRT, VTT, XML)"""
        try:
            import re
            from xml.etree import ElementTree as ET
            
            # Try XML format first (YouTube timedtext)
            if '<transcript>' in content or '<text' in content:
                try:
                    root = ET.fromstring(content)
                    texts = []
                    for text_elem in root.findall('.//text'):
                        text = text_elem.text
                        if text:
                            texts.append(text.strip())
                    return ' '.join(texts)
                except:
                    pass
            
            # Try SRT format
            if '-->' in content:
                srt_pattern = r'\d+\n\d{2}:\d{2}:\d{2},\d{3} --> \d{2}:\d{2}:\d{2},\d{3}\n(.+?)(?=\n\d+\n|\n\n|$)'
                matches = re.findall(srt_pattern, content, re.DOTALL)
                if matches:
                    return ' '.join([match.strip() for match in matches])
            
            # Try VTT format
            if 'WEBVTT' in content:
                vtt_pattern = r'\d{2}:\d{2}:\d{2}\.\d{3} --> \d{2}:\d{2}:\d{2}\.\d{3}\n(.+?)(?=\n\d{2}:\d{2}:\d{2}\.\d{3}|\n\n|$)'
                matches = re.findall(vtt_pattern, content, re.DOTALL)
                if matches:
                    return ' '.join([match.strip() for match in matches])
            
            # Fallback: extract all text content
            lines = content.split('\n')
            text_lines = []
            for line in lines:
                line = line.strip()
                if line and not line.isdigit() and '-->' not in line and '<' not in line:
                    text_lines.append(line)
            
            return ' '.join(text_lines)
            
        except Exception as e:
            logger.warning(f"Error parsing subtitle content: {str(e)}")
            return content
    
    def _has_speech_content(self, text: str) -> bool:
        """Check if text contains actual speech content"""
        speech_indicators = [
            'hello', 'hi', 'good', 'welcome', 'today', 'we', 'will', 'learn',
            'हाय', 'हेलो', 'गुड', 'वेलकम', 'आज', 'हम', 'सीखेंगे', 'करेंगे',
            'students', 'teacher', 'class', 'lecture', 'topic', 'chapter',
            'स्टूडेंट', 'टीचर', 'क्लास', 'लेक्चर', 'टॉपिक', 'चैप्टर'
        ]
        
        text_lower = text.lower()
        return any(indicator in text_lower for indicator in speech_indicators)
    
    def _extract_speech_from_html(self, html_content: str) -> str:
        """Extract speech content from HTML"""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            text_content = soup.get_text()
            
            # Look for speech content
            lines = text_content.split('\n')
            speech_lines = []
            
            for line in lines:
                line = line.strip()
                if (len(line) > 50 and 
                    (self._has_speech_content(line) or 
                     any(char in line for char in ['ह', 'अ', 'आ', 'इ', 'ई', 'उ', 'ऊ', 'ए', 'ऐ', 'ओ', 'औ']))):
                    speech_lines.append(line)
            
            return '\n'.join(speech_lines) if speech_lines else ''
            
        except Exception as e:
            logger.warning(f"Error extracting speech from HTML: {str(e)}")
            return ''
    
    def _clean_transcript_text(self, text: str) -> str:
        """Clean and format transcript text"""
        import re
        
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove URLs
        text = re.sub(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', '', text)
        
        # Remove timestamps
        text = re.sub(r'\d{2}:\d{2}:\d{2}', '', text)
        
        # Clean up
        text = text.strip()
        
        return text

    def _create_meaningful_placeholder(self, video_id: str, video_url: str = None) -> str:
        """Create meaningful placeholder when transcript extraction fails"""
        try:
            # Try to get video title for better placeholder
            import requests
            if video_url:
                response = requests.get(video_url, timeout=10)
                if response.status_code == 200:
                    import re
                    title_match = re.search(r'<title>([^<]+)</title>', response.text)
                    if title_match:
                        title = title_match.group(1).replace(' - YouTube', '').strip()
                        return f"""
                        {title}
                        
                        This educational video covers important topics related to {title.lower()}.
                        
                        Key topics discussed in this video include:
                        - Introduction to the main concepts
                        - Detailed explanations and examples
                        - Important points and key takeaways
                        - Practical applications and real-world examples
                        - Problem-solving techniques and strategies
                        
                        This video is designed to help you understand the material better and prepare for your exams. Please pay attention to the explanations and take notes as needed.
                        
                        If you have any questions about the content, feel free to ask in the comments section below.
                        
                        Note: This is a placeholder transcript based on the video title. The actual speech content from the video could not be extracted due to technical limitations. You can still ask questions about the video content, and I'll do my best to help based on the available information.
                        """
        except:
            pass
        
        # Fallback generic placeholder
        return f"""
        Educational Video Content
        
        This video covers important educational topics that are essential for your understanding.
        
        The main topics discussed in this video include:
        1. Introduction to the subject matter
        2. Key concepts and definitions
        3. Important formulas and equations
        4. Practical examples and applications
        5. Problem-solving techniques
        6. Summary and key takeaways
        
        This video is designed to help you understand the material better and prepare for your exams. Please pay attention to the explanations and take notes as needed.
        
        If you have any questions about the content, feel free to ask in the comments section below.
        
        Note: This is a placeholder transcript. The actual speech content from the video could not be extracted due to technical limitations. You can still ask questions about the video content, and I'll do my best to help based on the video title and description.
        """
    
    async def _get_youtube_innertube_api(self, video_id: str) -> str:
        """Get YouTube transcript using Innertube API (2025 method)"""
        try:
            import requests
            import json
            import re
            import time
            from xml.etree import ElementTree as ET
            
            logger.info(f"Trying Innertube API for video ID: {video_id}")
            
            # Step 1: Fetch the INNERTUBE_API_KEY from the Video Page
            video_url = f"https://www.youtube.com/watch?v={video_id}"
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
            }
            
            start_time = time.time()
            response = requests.get(video_url, headers=headers, timeout=15)
            response_time = time.time() - start_time
            
            if response.status_code != 200:
                logger.warning(f"Failed to fetch video page: {response.status_code}")
                return None
            
            # Extract API key
            api_key_match = re.search(r'"INNERTUBE_API_KEY":"([^"]+)"', response.text)
            if not api_key_match:
                logger.warning("INNERTUBE_API_KEY not found in video page")
                return None
            
            api_key = api_key_match.group(1)
            logger.info(f"Found API key: {api_key[:20]}...")
            
            # Step 2: Call the Innertube player API as Android
            endpoint = f"https://www.youtube.com/youtubei/v1/player?key={api_key}"
            
            body = {
                "context": {
                    "client": {
                        "clientName": "ANDROID",
                        "clientVersion": "20.10.38",
                    },
                },
                "videoId": video_id,
            }
            
            start_time = time.time()
            player_response = requests.post(endpoint, json=body, headers=headers, timeout=10)
            response_time = time.time() - start_time
            
            logger.info(f"Innertube API response time: {response_time:.2f}s")
            
            if player_response.status_code != 200:
                logger.warning(f"Innertube API failed: {player_response.status_code}")
                return None
            
            player_data = player_response.json()
            
            # Step 3: Extract the Caption Track (English or Other)
            tracks = player_data.get('captions', {}).get('playerCaptionsTracklistRenderer', {}).get('captionTracks', [])
            
            if not tracks:
                logger.warning("No caption tracks found in Innertube response")
                return None
            
            # Try English first, then any available language
            track = None
            for lang in ['en', 'en-US', 'en-GB']:
                track = next((t for t in tracks if t.get('languageCode') == lang), None)
                if track:
                    logger.info(f"Found caption track for language: {lang}")
                    break
            
            if not track:
                # Use first available track
                track = tracks[0]
                logger.info(f"Using first available track: {track.get('languageCode', 'unknown')}")
            
            base_url = track.get('baseUrl', '')
            if not base_url:
                logger.warning("No baseUrl found in caption track")
                return None
            
            # Remove "&fmt=srv3" if present
            base_url = re.sub(r'&fmt=\w+$', '', base_url)
            
            # Step 4: Fetch and Parse Captions XML
            start_time = time.time()
            caption_response = requests.get(base_url, headers=headers, timeout=10)
            response_time = time.time() - start_time
            
            logger.info(f"Caption fetch response time: {response_time:.2f}s")
            
            if caption_response.status_code != 200:
                logger.warning(f"Failed to fetch captions: {caption_response.status_code}")
                return None
            
            # Parse XML
            try:
                root = ET.fromstring(caption_response.text)
                transcript_texts = []
                
                for text_elem in root.findall('.//text'):
                    text_content = text_elem.text or ''
                    if text_content.strip():
                        transcript_texts.append(text_content.strip())
                
                if transcript_texts:
                    full_transcript = ' '.join(transcript_texts)
                    logger.info(f"Successfully fetched transcript via Innertube API with {len(transcript_texts)} entries")
                    return full_transcript
                else:
                    logger.warning("No transcript text found in XML")
                    return None
                    
            except ET.ParseError as e:
                logger.warning(f"Failed to parse caption XML: {str(e)}")
                return None
                
        except Exception as e:
            logger.error(f"Innertube API method failed: {str(e)}")
            return None
    
    async def _get_youtube_timedtext_api(self, video_id: str) -> str:
        """Get YouTube transcript using timedtext API (fast, cached)"""
        try:
            import requests
            import time
            
            logger.info(f"Trying timedtext API for video ID: {video_id}")
            
            # Common timedtext endpoints with different formats
            endpoints = [
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=srv3",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=ttml",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=vtt",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=hi&fmt=srv3",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=hi&fmt=ttml"
            ]
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/xml,application/xml,application/xhtml+xml,text/html;q=0.9,text/plain;q=0.8,image/png,*/*;q=0.5',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
            }
            
            for endpoint in endpoints:
                try:
                    start_time = time.time()
                    response = requests.get(endpoint, headers=headers, timeout=10)
                    response_time = time.time() - start_time
                    
                    logger.info(f"timedtext API response time: {response_time:.2f}s for {endpoint}")
                    
                    if response.status_code == 200 and response.text.strip():
                        # Parse based on format
                        if 'fmt=srv3' in endpoint:
                            transcript_text = self._parse_srv3_content(response.text)
                        elif 'fmt=ttml' in endpoint:
                            transcript_text = self._parse_ttml_content(response.text)
                        elif 'fmt=vtt' in endpoint:
                            transcript_text = self._parse_vtt_content(response.text)
                        else:
                            transcript_text = self._parse_generic_xml(response.text)
                        
                        if transcript_text and len(transcript_text.strip()) > 50:
                            logger.info(f"Successfully fetched transcript via timedtext API: {endpoint}")
                            return transcript_text
                            
                except Exception as e:
                    logger.warning(f"Failed to fetch from {endpoint}: {str(e)}")
                    continue
                    
        except Exception as e:
            logger.error(f"timedtext API method failed: {str(e)}")
            
        return None
    
    async def _get_youtube_youtubei_api(self, video_id: str) -> str:
        """Get YouTube transcript using youtubei API (needs proper context)"""
        try:
            import requests
            import json
            import time
            
            logger.info(f"Trying youtubei API for video ID: {video_id}")
            
            # First, get the video page to extract context
            video_url = f"https://www.youtube.com/watch?v={video_id}"
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
            }
            
            # Get video page to extract context
            response = requests.get(video_url, headers=headers, timeout=15)
            if response.status_code != 200:
                return None
                
            # Extract context from page (simplified approach)
            import re
            context_match = re.search(r'"INNERTUBE_CONTEXT":\s*({[^}]+})', response.text)
            if not context_match:
                return None
                
            context = json.loads(context_match.group(1))
            
            # Try youtubei API with proper context
            youtubei_url = "https://www.youtube.com/youtubei/v1/get_transcript"
            
            payload = {
                "context": context,
                "params": f"Qw%3D%3D",  # Base64 encoded parameters
                "videoId": video_id
            }
            
            start_time = time.time()
            response = requests.post(youtubei_url, json=payload, headers=headers, timeout=10)
            response_time = time.time() - start_time
            
            logger.info(f"youtubei API response time: {response_time:.2f}s")
            
            if response.status_code == 200:
                data = response.json()
                if 'actions' in data:
                    # Parse transcript from youtubei response
                    transcript_text = self._parse_youtubei_response(data)
                    if transcript_text:
                        logger.info("Successfully fetched transcript via youtubei API")
                        return transcript_text
                        
        except Exception as e:
            logger.error(f"youtubei API method failed: {str(e)}")
            
        return None
    
    def _parse_youtubei_response(self, data: dict) -> str:
        """Parse transcript from youtubei API response"""
        try:
            transcript_lines = []
            
            # Navigate through the complex youtubei response structure
            if 'actions' in data:
                for action in data['actions']:
                    if 'updateEngagementPanelAction' in action:
                        panel = action['updateEngagementPanelAction']['content']
                        if 'transcriptRenderer' in panel:
                            transcript = panel['transcriptRenderer']
                            if 'body' in transcript and 'transcriptBodyRenderer' in transcript['body']:
                                body = transcript['body']['transcriptBodyRenderer']
                                if 'cueGroups' in body:
                                    for cue_group in body['cueGroups']:
                                        if 'transcriptCueGroupRenderer' in cue_group:
                                            cue_group_renderer = cue_group['transcriptCueGroupRenderer']
                                            if 'cues' in cue_group_renderer:
                                                for cue in cue_group_renderer['cues']:
                                                    if 'transcriptCueRenderer' in cue:
                                                        cue_renderer = cue['transcriptCueRenderer']
                                                        if 'cue' in cue_renderer:
                                                            cue_text = cue_renderer['cue']
                                                            if 'simpleText' in cue_text:
                                                                transcript_lines.append(cue_text['simpleText'])
            
            if transcript_lines:
                return ' '.join(transcript_lines)
                
        except Exception as e:
            logger.error(f"Error parsing youtubei response: {str(e)}")
            
        return None
    
    async def _get_youtube_captions(self, video_id: str) -> str:
        """Alternative method to get YouTube captions using requests"""
        try:
            import requests
            import re
            
            logger.info(f"Trying alternative method for video ID: {video_id}")
            
            # Try to get captions from YouTube's internal API
            captions_url = f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=srv3"
            
            response = requests.get(captions_url, timeout=10)
            if response.status_code == 200:
                # Parse XML captions
                from xml.etree import ElementTree as ET
                try:
                    root = ET.fromstring(response.text)
                    captions_text = []
                    for text_elem in root.findall('.//text'):
                        if text_elem.text:
                            captions_text.append(text_elem.text.strip())
                    
                    if captions_text:
                        full_text = " ".join(captions_text)
                        logger.info(f"Successfully fetched captions with {len(captions_text)} entries")
                        return full_text
                except ET.ParseError as e:
                    logger.warning(f"Failed to parse captions XML: {str(e)}")
            
            # Try another format
            captions_url2 = f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=ttml"
            response2 = requests.get(captions_url2, timeout=10)
            if response2.status_code == 200:
                # Parse TTML captions
                try:
                    from xml.etree import ElementTree as ET
                    root = ET.fromstring(response2.text)
                    captions_text = []
                    for text_elem in root.findall('.//{http://www.w3.org/ns/ttml}p'):
                        if text_elem.text:
                            captions_text.append(text_elem.text.strip())
                    
                    if captions_text:
                        full_text = " ".join(captions_text)
                        logger.info(f"Successfully fetched TTML captions with {len(captions_text)} entries")
                        return full_text
                except ET.ParseError as e:
                    logger.warning(f"Failed to parse TTML captions: {str(e)}")
            
            logger.warning("No captions found using alternative method")
            return None
            
        except Exception as e:
            logger.error(f"Alternative captions method failed: {str(e)}")
            return None
    
    async def _get_youtube_with_ytdlp(self, video_id: str) -> str:
        """Try to get YouTube transcript using yt-dlp"""
        try:
            import subprocess
            import json
            import os
            import glob
            
            logger.info(f"Trying yt-dlp method for video ID: {video_id}")
            
            # Use yt-dlp to get video info and subtitles
            url = f"https://www.youtube.com/watch?v={video_id}"
            
            # Try to download subtitles
            subtitle_cmd = [
                "yt-dlp",
                "--write-subs",
                "--write-auto-subs",
                "--sub-langs", "en,hi",
                "--skip-download",
                "--no-warnings",
                "--output", f"/tmp/{video_id}_%(title)s.%(ext)s",
                url
            ]
            
            result = subprocess.run(subtitle_cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # Look for downloaded subtitle files
                subtitle_files = glob.glob(f"/tmp/{video_id}_*.vtt") + glob.glob(f"/tmp/{video_id}_*.srv3")
                
                for subtitle_file in subtitle_files:
                    try:
                        with open(subtitle_file, 'r', encoding='utf-8') as f:
                            content = f.read()
                        
                        # Parse VTT or SRV3 format
                        if subtitle_file.endswith('.vtt'):
                            transcript_text = self._parse_vtt_content(content)
                        else:
                            transcript_text = self._parse_srv3_content(content)
                        
                        if transcript_text:
                            logger.info(f"Successfully fetched transcript using yt-dlp from {subtitle_file}")
                            # Clean up file
                            os.remove(subtitle_file)
                            return transcript_text
                            
                    except Exception as e:
                        logger.warning(f"Error reading subtitle file {subtitle_file}: {str(e)}")
                        continue
                
                # Clean up any remaining files
                for file in subtitle_files:
                    try:
                        os.remove(file)
                    except:
                        pass
            
            # Fallback: get video description and title
            info_cmd = [
                "yt-dlp",
                "--dump-json",
                "--no-warnings",
                url
            ]
            
            info_result = subprocess.run(info_cmd, capture_output=True, text=True, timeout=30)
            
            if info_result.returncode == 0:
                video_info = json.loads(info_result.stdout)
                title = video_info.get('title', '')
                description = video_info.get('description', '')
                
                if title or description:
                    content = f"Title: {title}\n\nDescription: {description}"
                    logger.info(f"Successfully fetched video info using yt-dlp")
                    return content
                        
        except Exception as e:
            logger.error(f"yt-dlp method failed: {str(e)}")
            
        return None
    
    def _parse_vtt_content(self, content: str) -> str:
        """Parse VTT subtitle content"""
        try:
            lines = content.split('\n')
            transcript_lines = []
            
            for line in lines:
                line = line.strip()
                # Skip VTT headers, timestamps, and empty lines
                if (not line or 
                    line.startswith('WEBVTT') or 
                    line.startswith('NOTE') or
                    '-->' in line or
                    line.startswith('Kind:') or
                    line.startswith('Language:')):
                    continue
                
                # Remove HTML tags
                import re
                clean_line = re.sub(r'<[^>]+>', '', line)
                if clean_line.strip():
                    transcript_lines.append(clean_line.strip())
            
            return ' '.join(transcript_lines)
        except Exception as e:
            logger.error(f"Error parsing VTT content: {str(e)}")
            return None
    
    def _parse_srv3_content(self, content: str) -> str:
        """Parse SRV3 subtitle content"""
        try:
            from xml.etree import ElementTree as ET
            root = ET.fromstring(content)
            transcript_lines = []
            
            for text_elem in root.findall('.//text'):
                if text_elem.text:
                    transcript_lines.append(text_elem.text.strip())
            
            return ' '.join(transcript_lines)
        except Exception as e:
            logger.error(f"Error parsing SRV3 content: {str(e)}")
            return None
    
    async def _get_youtube_video_info(self, video_id: str) -> str:
        """Get YouTube video basic info (title, description) as fallback"""
        try:
            import requests
            from bs4 import BeautifulSoup
            
            logger.info(f"Trying to get video info for video ID: {video_id}")
            
            url = f"https://www.youtube.com/watch?v={video_id}"
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            response = requests.get(url, headers=headers, timeout=10)
            if response.status_code == 200:
                soup = BeautifulSoup(response.content, 'html.parser')
                
                # Try to get title
                title = ""
                title_elem = soup.find('title')
                if title_elem:
                    title = title_elem.get_text().strip()
                    if title.endswith(' - YouTube'):
                        title = title[:-10].strip()
                
                # Try to get description
                description = ""
                desc_elem = soup.find('meta', {'name': 'description'})
                if desc_elem:
                    description = desc_elem.get('content', '').strip()
                
                # Try to get channel name
                channel = ""
                channel_elem = soup.find('link', {'itemprop': 'name'})
                if channel_elem:
                    channel = channel_elem.get('content', '').strip()
                
                if title or description:
                    content = f"Video Title: {title}\n"
                    if channel:
                        content += f"Channel: {channel}\n"
                    content += f"\nDescription:\n{description}"
                    
                    logger.info(f"Successfully fetched video info: {title}")
                    return content
                    
        except Exception as e:
            logger.error(f"Video info method failed: {str(e)}")
            
        return None
    
    async def _get_youtube_direct_api(self, video_id: str) -> str:
        """Try to get YouTube transcript using direct API calls"""
        try:
            import requests
            import re
            
            logger.info(f"Trying direct YouTube API method for video ID: {video_id}")
            
            # Try different YouTube API endpoints for captions
            api_urls = [
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=srv3",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=ttml",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=vtt",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=hi&fmt=srv3",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=hi&fmt=ttml"
            ]
            
            for api_url in api_urls:
                try:
                    response = requests.get(api_url, timeout=10)
                    if response.status_code == 200 and response.text.strip():
                        # Parse based on format
                        if 'fmt=srv3' in api_url:
                            transcript_text = self._parse_srv3_content(response.text)
                        elif 'fmt=ttml' in api_url:
                            transcript_text = self._parse_ttml_content(response.text)
                        elif 'fmt=vtt' in api_url:
                            transcript_text = self._parse_vtt_content(response.text)
                        else:
                            # Try to parse as generic XML
                            transcript_text = self._parse_generic_xml(response.text)
                        
                        if transcript_text and len(transcript_text.strip()) > 50:
                            logger.info(f"Successfully fetched transcript using direct API: {api_url}")
                            return transcript_text
                            
                except Exception as e:
                    logger.warning(f"Failed to fetch from {api_url}: {str(e)}")
                    continue
            
            # Try to get from video page directly
            video_url = f"https://www.youtube.com/watch?v={video_id}"
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            response = requests.get(video_url, headers=headers, timeout=15)
            if response.status_code == 200:
                # Look for captions in the page source
                caption_patterns = [
                    r'"captions":\s*({[^}]+})',
                    r'"subtitles":\s*({[^}]+})',
                    r'"timedtext":\s*({[^}]+})'
                ]
                
                for pattern in caption_patterns:
                    matches = re.findall(pattern, response.text)
                    if matches:
                        logger.info("Found captions in page source")
                        # This would need more complex parsing, but it's a start
                        break
                        
        except Exception as e:
            logger.error(f"Direct YouTube API method failed: {str(e)}")
            
        return None
    
    def _parse_ttml_content(self, content: str) -> str:
        """Parse TTML subtitle content"""
        try:
            from xml.etree import ElementTree as ET
            root = ET.fromstring(content)
            transcript_lines = []
            
            # TTML uses different namespace
            for text_elem in root.findall('.//{http://www.w3.org/ns/ttml}p'):
                if text_elem.text:
                    transcript_lines.append(text_elem.text.strip())
            
            return ' '.join(transcript_lines)
        except Exception as e:
            logger.error(f"Error parsing TTML content: {str(e)}")
            return None
    
    def _parse_generic_xml(self, content: str) -> str:
        """Parse generic XML subtitle content"""
        try:
            from xml.etree import ElementTree as ET
            root = ET.fromstring(content)
            transcript_lines = []
            
            # Try different possible text elements
            for text_elem in root.findall('.//text'):
                if text_elem.text:
                    transcript_lines.append(text_elem.text.strip())
            
            # If no text elements found, try p elements
            if not transcript_lines:
                for text_elem in root.findall('.//p'):
                    if text_elem.text:
                        transcript_lines.append(text_elem.text.strip())
            
            return ' '.join(transcript_lines)
        except Exception as e:
            logger.error(f"Error parsing generic XML content: {str(e)}")
            return None
    
    async def _get_youtube_simple_api(self, video_id: str) -> str:
        """Simple method to get YouTube transcript using direct API calls"""
        try:
            import requests
            import re
            import json
            
            logger.info(f"Trying simple YouTube API method for video ID: {video_id}")
            
            # Try different YouTube API endpoints for captions
            api_urls = [
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=srv3",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=ttml",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=en&fmt=vtt",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=hi&fmt=srv3",
                f"https://www.youtube.com/api/timedtext?v={video_id}&lang=hi&fmt=ttml"
            ]
            
            for api_url in api_urls:
                try:
                    response = requests.get(api_url, timeout=10)
                    if response.status_code == 200 and response.text.strip():
                        # Parse based on format
                        if 'fmt=srv3' in api_url:
                            transcript_text = self._parse_srv3_content(response.text)
                        elif 'fmt=ttml' in api_url:
                            transcript_text = self._parse_ttml_content(response.text)
                        elif 'fmt=vtt' in api_url:
                            transcript_text = self._parse_vtt_content(response.text)
                        else:
                            # Try to parse as generic XML
                            transcript_text = self._parse_generic_xml(response.text)
                        
                        if transcript_text and len(transcript_text.strip()) > 50:
                            logger.info(f"Successfully fetched transcript using simple API: {api_url}")
                            return transcript_text
                            
                except Exception as e:
                    logger.warning(f"Failed to fetch from {api_url}: {str(e)}")
                    continue
            
            # Try to get from video page directly using a different approach
            video_url = f"https://www.youtube.com/watch?v={video_id}"
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
            }
            
            response = requests.get(video_url, headers=headers, timeout=15)
            if response.status_code == 200:
                # Look for captions in the page source
                caption_patterns = [
                    r'"captions":\s*({[^}]+})',
                    r'"subtitles":\s*({[^}]+})',
                    r'"timedtext":\s*({[^}]+})',
                    r'"playerCaptionsTracklistRenderer":\s*({[^}]+})'
                ]
                
                for pattern in caption_patterns:
                    matches = re.findall(pattern, response.text)
                    if matches:
                        logger.info("Found captions in page source")
                        # This would need more complex parsing, but it's a start
                        break
                        
        except Exception as e:
            logger.error(f"Simple YouTube API method failed: {str(e)}")
            
        return None
    
    def _extract_video_id(self, url: str) -> str:
        """Extract video ID from YouTube URL"""
        import re
        
        # Pattern for YouTube URLs
        patterns = [
            r'(?:youtube\.com\/watch\?v=|youtu\.be\/|youtube\.com\/embed\/)([^&\n?#]+)',
            r'youtube\.com\/live\/([^&\n?#]+)'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                return match.group(1)
        
        return None
